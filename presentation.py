from pptx import Presentation



#sample input, dictionary of a dictionary - index i of input = {header:{subheader:bullet_points}} or
# input = {
#     'header1': {
#         'subheader_1': ['bullet_point_1', 'bullet_point_2'],
#         'subheader_2': ['bullet_point_1']
#     },
#     'header2': {
#         'subheader_3': ['bullet_point_1', 'bullet_point_2', 'bullet_point_3'],
#         'subheader_4': ['bullet_point_1']
#     }
# }
def add_bullet_slide(prs, title, input):
    bullet_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    
    shapes.title.text = title

    body_shape = shapes.placeholders[1]
    tf = body_shape.text_frame

    # Use the first paragraph for the first header
    first = True
    for header, content in input.items():
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.text = header 
        p.level = 0
        
        if not isinstance(content, set):
            for subheader, bullet_points in content.items():
                p = tf.add_paragraph()
                p.text = subheader
                p.level = 1

                if bullet_points:  # Check if the list is not empty
                    for bullet in bullet_points:
                        p = tf.add_paragraph()
                        p.text = bullet
                        p.level = 2
        else:
            for subheader in content:
                p = tf.add_paragraph()
                p.text = subheader
                p.level = 1


prs = Presentation()

input = {
    "Structure and Composition:" : {
        "Congress" : ["the House of Representatives", "The Senate"]
    },
    "The house" : {
        "Members"
    }
}

# Structure and Composition:
# Congress consists of two chambers: the House of Representatives and the Senate
# House has 435 voting members, apportioned by state population
# Senate has 100 members, 2 from each state
# Members must meet age, citizenship, and residency requirements

add_bullet_slide(prs, "Hello", input)

prs.save('test.pptx')